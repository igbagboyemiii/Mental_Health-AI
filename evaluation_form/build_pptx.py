"""
build_pptx.py  —  MindGuard Thesis Defence Presentation Generator
Run: python evaluation_form/build_pptx.py
Output: evaluation_form/MindGuard_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x0D, 0x25)   # dark navy
C_CARD      = RGBColor(0x1A, 0x1A, 0x3A)   # card bg
C_PRIMARY   = RGBColor(0x6C, 0x63, 0xFF)   # indigo
C_SECONDARY = RGBColor(0xA7, 0x8B, 0xFA)   # violet
C_ACCENT    = RGBColor(0x34, 0xD3, 0x99)   # emerald
C_DANGER    = RGBColor(0xF8, 0x71, 0x71)   # red
C_WARNING   = RGBColor(0xFB, 0xBF, 0x24)   # amber
C_TEXT      = RGBColor(0xE2, 0xE8, 0xF0)   # light grey
C_MUTED     = RGBColor(0x94, 0xA3, 0xB8)   # slate
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_size=14, bold=False, color=C_TEXT,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Calibri"
    return txBox

def add_para_to_tf(tf, text, font_size=12, bold=False,
                   color=C_TEXT, align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Calibri"
    return p

def add_label_badge(slide, text, l, t):
    """Small purple tag label."""
    w, h = Inches(1.5), Inches(0.28)
    r = add_rect(slide, l, t, w, h, C_PRIMARY)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = C_WHITE
    run.font.name = "Calibri"

def slide_header(slide, section_tag, title_parts):
    """Draw the section badge + title + horizontal rule."""
    add_label_badge(slide, section_tag, Inches(0.5), Inches(0.45))
    # Title
    txBox = slide.shapes.add_textbox(Inches(2.2), Inches(0.38), Inches(9.8), Inches(0.55))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    for part in title_parts:
        run = p.add_run()
        run.text = part[0]
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = part[1]
        run.font.name = "Calibri"
    # Horizontal rule (thin rect)
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(12.3), Pt(1.5), C_PRIMARY)

def card_box(slide, l, t, w, h, icon, title, body, icon_size=22, title_size=12, body_size=10.5):
    add_rect(slide, l, t, w, h, C_CARD, C_PRIMARY, Pt(0.5))
    # icon
    add_textbox(slide, icon, l+Inches(0.15), t+Inches(0.12), Inches(0.5), Inches(0.4),
                font_size=icon_size, color=C_WHITE)
    # title
    add_textbox(slide, title, l+Inches(0.15), t+Inches(0.52), w-Inches(0.3), Inches(0.28),
                font_size=title_size, bold=True, color=C_WHITE)
    # body
    tb = slide.shapes.add_textbox(l+Inches(0.15), t+Inches(0.82), w-Inches(0.3), h-Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = body
    run.font.size = Pt(body_size)
    run.font.color.rgb = C_MUTED
    run.font.name = "Calibri"

def bullet_item(slide, num_text, body, cite, l, t, w, h,
                num_color=C_PRIMARY, border_color=C_PRIMARY):
    add_rect(slide, l, t, w, h, C_CARD, border_color, Pt(0.5))
    # left accent stripe
    add_rect(slide, l, t, Pt(4), h, border_color)
    # number circle
    add_rect(slide, l+Inches(0.15), t+h/2-Inches(0.18), Inches(0.36), Inches(0.36),
             RGBColor(0x1E, 0x1E, 0x45), num_color, Pt(0.5))
    add_textbox(slide, num_text, l+Inches(0.15), t+h/2-Inches(0.18),
                Inches(0.36), Inches(0.36), font_size=9, bold=True,
                color=num_color, align=PP_ALIGN.CENTER)
    # body text
    tb = slide.shapes.add_textbox(l+Inches(0.6), t+Inches(0.1), w-Inches(0.75), h-Inches(0.32))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(11)
    run.font.color.rgb = C_TEXT
    run.font.name = "Calibri"
    if cite:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = cite
        run2.font.size = Pt(9)
        run2.font.color.rgb = C_PRIMARY
        run2.font.italic = True
        run2.font.name = "Calibri"

def metric_box(slide, l, t, w, h, value, label):
    add_rect(slide, l, t, w, h, C_CARD, C_PRIMARY, Pt(0.5))
    add_textbox(slide, value, l, t+Inches(0.18), w, Inches(0.55),
                font_size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, label.upper(), l, t+Inches(0.75), w, Inches(0.28),
                font_size=8.5, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

def table_row(slide, cols, l, t, w, h, is_header=False, highlight=False):
    col_w = w / len(cols)
    bg = RGBColor(0x1A, 0x1A, 0x45) if is_header else (C_CARD if not highlight else RGBColor(0x18, 0x18, 0x3A))
    for i, text in enumerate(cols):
        add_rect(slide, l + col_w*i, t, col_w, h, bg, C_PRIMARY, Pt(0.3))
        color = C_SECONDARY if is_header else C_TEXT
        add_textbox(slide, text, l+col_w*i+Inches(0.08), t+Inches(0.06),
                    col_w-Inches(0.16), h-Inches(0.12),
                    font_size=9 if not is_header else 8,
                    bold=is_header, color=color, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)

# gradient overlay rectangles for visual depth
add_rect(sl, 0, 0, SLIDE_W, Inches(3.5), RGBColor(0x10, 0x10, 0x28))
add_rect(sl, 0, Inches(3.5), SLIDE_W, Inches(4.0), RGBColor(0x0A, 0x0A, 0x1E))

# decorative circle accents
add_rect(sl, Inches(10.5), Inches(-1), Inches(4), Inches(4), RGBColor(0x15,0x12,0x40))

# Badge
badge = add_rect(sl, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.3), RGBColor(0x1A,0x14,0x45), C_PRIMARY, Pt(0.5))
add_textbox(sl, "FINAL YEAR BSc PROJECT DEFENCE · 2026",
            Inches(4.5), Inches(1.22), Inches(4.3), Inches(0.3),
            font_size=8, bold=True, color=C_SECONDARY, align=PP_ALIGN.CENTER)

# Emoji / icon
add_textbox(sl, "🧠🛡️", Inches(5.8), Inches(1.62), Inches(2), Inches(0.7),
            font_size=30, align=PP_ALIGN.CENTER, color=C_WHITE)

# Title
add_textbox(sl, "MindGuard", Inches(1.5), Inches(2.45), Inches(10.3), Inches(1.0),
            font_size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(sl,
    "Development of a Browser Extension for Detecting Depressive Language\nin Adolescents in Online Communication Using a RAG-Driven LLM",
    Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
    font_size=13, color=C_MUTED, align=PP_ALIGN.CENTER)

# Divider
add_rect(sl, Inches(5.9), Inches(4.38), Inches(1.5), Pt(3), C_PRIMARY)

# Author
add_textbox(sl, "Adelekan Igbagboyemi Mary",
            Inches(1.5), Inches(4.55), Inches(10.3), Inches(0.38),
            font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, "Computer Science  ·  Final Year  ·  2025/2026",
            Inches(1.5), Inches(4.96), Inches(10.3), Inches(0.3),
            font_size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

# Slide number
add_textbox(sl, "01", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCTION / BACKGROUND
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Introduction", [("Background to the ", C_WHITE), ("Field", C_PRIMARY)])
add_textbox(sl, "02", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

cw = Inches(3.95)
ch = Inches(2.4)
ct = Inches(1.25)
card_box(sl, Inches(0.5),  ct, cw, ch, "🌍",
         "Global Mental Health Crisis",
         "Depression affects 280 million people worldwide. Adolescents are disproportionately impacted — 1 in 7 aged 10–19 experience a mental disorder, yet less than 20% receive adequate care.")

card_box(sl, Inches(4.68), ct, cw, ch, "📱",
         "Adolescents & Online Communication",
         "Adolescents spend 7+ hours/day online. They frequently express emotional distress through written online text — often their primary and most candid outlet for self-expression.")

card_box(sl, Inches(8.86), ct, cw, ch, "⚠️",
         "The Intervention Gap",
         "Existing tools rely on active user engagement — journaling apps, chatbots — but at-risk adolescents disengage when most distressed. Early detection requires passive, non-intrusive monitoring.")

# Quote
add_rect(sl, Inches(0.5), Inches(3.82), Inches(12.3), Inches(0.85), C_CARD, C_PRIMARY, Pt(0.5))
add_rect(sl, Inches(0.5), Inches(3.82), Pt(4), Inches(0.85), C_PRIMARY)
add_textbox(sl,
    '"The digital footprint left by adolescents during online communication represents an untapped early-warning signal for identifying depressive risk before it escalates to crisis."',
    Inches(0.68), Inches(3.9), Inches(12.0), Inches(0.72),
    font_size=10.5, italic=True, color=C_TEXT)

# Flow diagram
flow_items = ["👤 Writes Online", "→", "🔍 Passive Capture", "→", "🤖 AI Analysis", "→", "🔔 Guardian Alert", "→", "💚 Intervention"]
flow_colors = [C_CARD, C_BG, C_CARD, C_BG, C_CARD, C_BG, C_CARD, C_BG, C_CARD]
fx = Inches(0.5)
for i, item in enumerate(flow_items):
    if "→" in item:
        add_textbox(sl, "→", fx, Inches(4.88), Inches(0.3), Inches(0.35),
                    font_size=14, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
        fx += Inches(0.32)
    else:
        bw = Inches(1.58)
        add_rect(sl, fx, Inches(4.82), bw, Inches(0.45), C_CARD, C_PRIMARY, Pt(0.5))
        add_textbox(sl, item, fx, Inches(4.88), bw, Inches(0.35),
                    font_size=9, color=C_TEXT, align=PP_ALIGN.CENTER)
        fx += bw + Inches(0.02)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — LITERATURE REVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Literature Review", [("Key Works ", C_WHITE), ("Reviewed", C_PRIMARY)])
add_textbox(sl, "03", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

headers = ["Author(s) & Year", "Problem Considered", "Method Used", "Gap Identified", "Addressed by MindGuard?"]
col_ws  = [Inches(2.0), Inches(2.6), Inches(2.6), Inches(2.6), Inches(2.5)]

tl = Inches(0.5)
tt = Inches(1.22)
th = Inches(0.4)

# Header row
for i, h in enumerate(headers):
    xl = tl + sum(col_ws[:i])
    add_rect(sl, xl, tt, col_ws[i], th, RGBColor(0x1A,0x1A,0x45), C_PRIMARY, Pt(0.3))
    add_textbox(sl, h, xl+Inches(0.06), tt+Inches(0.06),
                col_ws[i]-Inches(0.12), th-Inches(0.1),
                font_size=8.5, bold=True, color=C_SECONDARY)

rows = [
    ["Q. Guo et al. (2024)",    "Suicidal ideation detection in social media",   "BERT-based classification on Reddit",         "Not age-targeted; no guardian alerts",           "✔ Age-targeted + Alerts"],
    ["Ilapaka & Ghosh (2025)",  "AI chatbot for depression support",              "FAISS + RAG + active dialogue",               "Requires active user initiation",                "✔ Passive monitoring"],
    ["Zhang et al. — SouLLMate (2023)", "LLM emotional support for mental health","Fine-tuned LLM; overall F1 = 0.68",          "No guardian escalation; low F1",                "✔ Guardian dashboard + 94.1% High-Risk F1"],
    ["Fitzpatrick et al. — Woebot (2022)","Automated CBT chatbot for depression", "Rule-based NLP; brief daily conversations",  "Requires active engagement; no passive mode",   "✔ Zero user effort required"],
    ["Hollis et al. (2021)",    "Digital mental health for young people: review","Systematic review; evidence synthesis",        "Gap in passive real-time browser tools for teens","✔ Browser extension solution"],
]

row_h = Inches(0.95)
for ri, row in enumerate(rows):
    ry = tt + th + ri * row_h
    bg = RGBColor(0x18,0x18,0x3A) if ri % 2 == 0 else C_CARD
    for ci, cell in enumerate(row):
        xl = tl + sum(col_ws[:ci])
        add_rect(sl, xl, ry, col_ws[ci], row_h, bg, C_PRIMARY, Pt(0.2))
        color = C_ACCENT if ci == 4 else C_TEXT
        add_textbox(sl, cell, xl+Inches(0.07), ry+Inches(0.06),
                    col_ws[ci]-Inches(0.14), row_h-Inches(0.1),
                    font_size=9, color=color, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — STATEMENT OF PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Statement of Problem", [("The ", C_WHITE), ("Problem ", C_PRIMARY), ("We Solve", C_WHITE)])
add_textbox(sl, "04", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

problems = [
    ("1", "Active-Engagement Dependency",
     "Existing digital mental health tools work only if the user chooses to engage. Adolescents experiencing depressive episodes are the least likely to do so — the very moment help is needed most.",
     "(Q. Guo et al., 2024; Hollis et al., 2021)"),
    ("2", "Passive Expression Goes Undetected",
     "Adolescents express distress passively through everyday online writing. No browser-based tool exists to passively capture and analyse this real-time signal without user effort.",
     "(O'Reilly et al., 2023; Torous et al., 2021)"),
    ("3", "No Guardian Escalation Framework",
     "State-of-the-art detection models lack an integrated guardian notification system. Detection alone does not bridge the gap between identifying risk and real-world intervention.",
     "(Zhang et al., 2023)"),
    ("4", "Not Designed for Adolescents",
     "Current tools ignore the informal, colloquial, and coded language unique to teenage online communication, reducing detection reliability for the most vulnerable demographic.",
     "(Thabrew et al., 2022)"),
]

bt = Inches(1.3)
bh = Inches(1.28)
for i, (num, title, body, cite) in enumerate(problems):
    y = bt + i * (bh + Inches(0.06))
    add_rect(sl, Inches(0.5), y, Inches(12.3), bh, C_CARD, C_PRIMARY, Pt(0.4))
    add_rect(sl, Inches(0.5), y, Pt(4), bh, C_PRIMARY)
    # Number badge
    add_rect(sl, Inches(0.65), y+Inches(0.38), Inches(0.35), Inches(0.35),
             RGBColor(0x1E,0x1E,0x45), C_PRIMARY, Pt(0.5))
    add_textbox(sl, num, Inches(0.65), y+Inches(0.38), Inches(0.35), Inches(0.35),
                font_size=10, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    # Title
    add_textbox(sl, title, Inches(1.15), y+Inches(0.1), Inches(11.4), Inches(0.3),
                font_size=12, bold=True, color=C_WHITE)
    # Body
    add_textbox(sl, body, Inches(1.15), y+Inches(0.42), Inches(11.0), Inches(0.52),
                font_size=10.5, color=C_TEXT, wrap=True)
    # Citation
    add_textbox(sl, cite, Inches(1.15), y+Inches(0.95), Inches(11.0), Inches(0.26),
                font_size=9, color=C_PRIMARY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AIM
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Aim", [("What We Set Out to ", C_WHITE), ("Achieve", C_PRIMARY)])
add_textbox(sl, "05", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

# Central aim box
add_rect(sl, Inches(1.5), Inches(1.5), Inches(10.3), Inches(3.2),
         C_CARD, C_PRIMARY, Pt(1.5))
add_textbox(sl, "🎯", Inches(6.1), Inches(1.65), Inches(1.1), Inches(0.7),
            font_size=32, align=PP_ALIGN.CENTER, color=C_WHITE)
add_textbox(sl,
    "To design, develop, and evaluate a browser extension that passively monitors online text for depressive language patterns in adolescents, using a RAG-driven Large Language Model, and automatically notifies designated guardians when high-risk indicators are detected.",
    Inches(1.9), Inches(2.45), Inches(9.5), Inches(1.8),
    font_size=14.5, color=C_TEXT, align=PP_ALIGN.CENTER, wrap=True)

# Chips
chips = ["🔍  Passive Monitoring", "🤖  RAG + LLM", "🛡️  Risk Classification", "🔔  Guardian Alerts"]
cw_chip = Inches(2.6)
start_x = Inches(0.5) + (SLIDE_W - cw_chip*4 - Inches(0.2)*3) / 2
for i, chip in enumerate(chips):
    cx = start_x + i*(cw_chip + Inches(0.2))
    add_rect(sl, cx, Inches(4.92), cw_chip, Inches(0.42),
             RGBColor(0x1A,0x14,0x45), C_PRIMARY, Pt(0.5))
    add_textbox(sl, chip, cx, Inches(4.95), cw_chip, Inches(0.36),
                font_size=11, bold=True, color=C_SECONDARY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — OBJECTIVES
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Objectives", [("Steps Taken to ", C_WHITE), ("Achieve the Aim", C_PRIMARY)])
add_textbox(sl, "06", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

objectives = [
    ("I",   "Collect & Pre-process Dataset",
     "Collect and pre-process a suitable dataset of depressive and non-depressive text samples representative of adolescent online communication (Dreaddit dataset — 2,838 samples)."),
    ("II",  "Design & Develop Risk Model",
     "Design and develop a trained risk classification model and evaluate the accuracy and reliability of the model's risk predictions using a hybrid NLP + RAG pipeline."),
    ("III", "Implement System Architecture",
     "Implement the model within a browser extension architecture, incorporating a secure guardian notification framework and consent management system."),
    ("IV",  "Evaluate Usability",
     "Evaluate the usability of the system interface using the System Usability Scale (SUS) with 24 participants from the target demographic."),
]

ot = Inches(1.28)
oh = Inches(1.3)
for i, (num, title, body) in enumerate(objectives):
    y = ot + i*(oh + Inches(0.06))
    add_rect(sl, Inches(0.5), y, Inches(12.3), oh, C_CARD, C_PRIMARY, Pt(0.4))
    add_rect(sl, Inches(0.5), y, Pt(4), oh, C_PRIMARY)
    add_rect(sl, Inches(0.65), y+Inches(0.35), Inches(0.42), Inches(0.42),
             RGBColor(0x1E,0x1E,0x45), C_PRIMARY, Pt(0.5))
    add_textbox(sl, num, Inches(0.65), y+Inches(0.35), Inches(0.42), Inches(0.42),
                font_size=10, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    add_textbox(sl, title, Inches(1.22), y+Inches(0.12), Inches(11.2), Inches(0.28),
                font_size=12, bold=True, color=C_WHITE)
    add_textbox(sl, body,  Inches(1.22), y+Inches(0.42), Inches(11.2), Inches(0.72),
                font_size=10.5, color=C_TEXT, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Methodology", [("How Each Objective Was ", C_WHITE), ("Achieved", C_ACCENT)])
add_textbox(sl, "07", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

# Column headers
add_rect(sl, Inches(0.5),  Inches(1.22), Inches(5.8), Inches(0.3), RGBColor(0x1A,0x14,0x45), C_PRIMARY, Pt(0.3))
add_textbox(sl, "OBJECTIVE", Inches(0.5), Inches(1.24), Inches(5.8), Inches(0.26),
            font_size=8, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(6.5),  Inches(1.22), Inches(6.3), Inches(0.3), RGBColor(0x0D,0x25,0x1A), C_ACCENT, Pt(0.3))
add_textbox(sl, "METHOD USED", Inches(6.5), Inches(1.24), Inches(6.3), Inches(0.26),
            font_size=8, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

meth_rows = [
    ("Obj I — Dataset",
     "Dreaddit Reddit dataset · 2,838 samples · Colloquial adolescent-style text",
     "Text cleaning, stopword removal, lemmatisation. FAISS embedding index built using sentence-transformer embeddings (all-MiniLM-L6-v2)."),
    ("Obj II — Risk Model",
     "3-class Risk Engine: Low / Moderate / High. Hybrid NLP scoring + RAG retrieval.",
     "RAG Pipeline: FAISS vector search retrieves semantically similar cases → Gemini 1.5 Flash reasons over context → produces risk level + explanation. 14-day temporal window (DSM-5 aligned)."),
    ("Obj III — System",
     "Browser extension (JS), FastAPI backend, Guardian dashboard (HTML/JS).",
     "Design Science Research (DSR). Passive text interception → encrypted API call → risk classification → conditional guardian email alert with consent verification."),
    ("Obj IV — Evaluation",
     "SUS questionnaire + model benchmarking on 42-sample held-out test set.",
     "Quantitative: Accuracy, Precision, Recall, F1, Confusion Matrix. Qualitative: 10-item SUS with 24 participants on live extension + guardian dashboard."),
]

mt = Inches(1.62)
mh = Inches(1.2)
for i, (obj_label, obj_body, meth_body) in enumerate(meth_rows):
    y = mt + i*(mh+Inches(0.05))
    # Objective box (left)
    add_rect(sl, Inches(0.5), y, Inches(5.8), mh,
             RGBColor(0x14,0x14,0x35), C_PRIMARY, Pt(0.4))
    add_textbox(sl, obj_label, Inches(0.65), y+Inches(0.08), Inches(5.5), Inches(0.24),
                font_size=9, bold=True, color=C_PRIMARY)
    add_textbox(sl, obj_body, Inches(0.65), y+Inches(0.34), Inches(5.4), mh-Inches(0.42),
                font_size=10, color=C_TEXT, wrap=True)
    # Arrow
    add_textbox(sl, "→", Inches(6.15), y + mh/2 - Inches(0.18), Inches(0.3), Inches(0.36),
                font_size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    # Method box (right)
    add_rect(sl, Inches(6.5), y, Inches(6.3), mh,
             RGBColor(0x0D,0x20,0x18), C_ACCENT, Pt(0.4))
    add_textbox(sl, meth_body, Inches(6.65), y+Inches(0.1), Inches(6.0), mh-Inches(0.18),
                font_size=10, color=C_TEXT, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS: Obj I & II (Classification Metrics)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Results", [("Objective I & II — ", C_WHITE), ("Model Performance", C_ACCENT)])
add_textbox(sl, "08", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

# Metric boxes row
mbox_w = Inches(2.85)
mbox_h = Inches(0.95)
metrics = [("88.1%","Accuracy"),("87.4%","Precision"),("88.1%","Recall"),("87.5%","F1-Score")]
for i,(val,lbl) in enumerate(metrics):
    metric_box(sl, Inches(0.5)+i*(mbox_w+Inches(0.1)), Inches(1.22), mbox_w, mbox_h, val, lbl)

# Per-class table
th2 = Inches(2.32)
col_ws2 = [Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(5.1)]
hdr2 = ["Class", "Precision", "Recall", "F1-Score", "Clinical Significance"]
row_h2 = Inches(0.38)
for ci, h in enumerate(hdr2):
    xl = Inches(0.5) + sum(col_ws2[:ci])
    add_rect(sl, xl, th2, col_ws2[ci], row_h2, RGBColor(0x1A,0x1A,0x45), C_PRIMARY, Pt(0.3))
    add_textbox(sl, h, xl+Inches(0.06), th2+Inches(0.06), col_ws2[ci]-Inches(0.12), row_h2-Inches(0.1),
                font_size=8.5, bold=True, color=C_SECONDARY)

table_data = [
    (["Low Risk","73.3%","84.6%","78.5%",
      "26.7% false-positive rate — identified target for future improvement to reduce guardian fatigue"],
     C_ACCENT, False),
    (["Moderate Risk","90.0%","82.6%","86.1%",
      "Strong identification of borderline cases — supports preventative intervention before escalation"],
     C_WARNING, True),
    (["High Risk","94.4%","94.1%","94.1%",
      "✔ MEETS ≥90% clinical safety target. High Recall is paramount: missing a high-risk user has severe consequences."],
     C_DANGER, False),
]

for ri,(row,label_color,highlight) in enumerate(table_data):
    ry = th2 + row_h2 + ri*row_h2
    bg = RGBColor(0x18,0x18,0x3A) if highlight else C_CARD
    for ci, cell in enumerate(row):
        xl = Inches(0.5) + sum(col_ws2[:ci])
        add_rect(sl, xl, ry, col_ws2[ci], row_h2, bg, C_PRIMARY, Pt(0.2))
        color = label_color if ci==0 else (C_ACCENT if ci==4 and ri==2 else C_TEXT)
        font_b = (ci == 0)
        add_textbox(sl, cell, xl+Inches(0.06), ry+Inches(0.04),
                    col_ws2[ci]-Inches(0.12), row_h2-Inches(0.06),
                    font_size=9, bold=font_b, color=color, wrap=True)

# Footer info
add_textbox(sl,
    "📊 Test samples: 42 hand-curated high-variance instances   ·   🔬 Dataset: Dreaddit (2,838 total samples)   ·   ⏱ 14-day temporal window (DSM-5 aligned)",
    Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.32),
    font_size=9, color=C_MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — RESULTS: Obj III & IV (System + SUS)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Results", [("Objective III & IV — ", C_WHITE), ("System & Usability", C_ACCENT)])
add_textbox(sl, "09", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

# Left: System delivered
add_textbox(sl, "OBJECTIVE III — SYSTEM DELIVERED", Inches(0.5), Inches(1.22), Inches(5.8), Inches(0.26),
            font_size=8, bold=True, color=C_PRIMARY)
sys_items = [
    ("🔌", "Browser Extension", "Passively captures typed text. Encrypts & transmits to API. Supports Chrome."),
    ("🤖", "RAG Risk Engine",   "FAISS retrieval + Gemini LLM classifies text into Low / Moderate / High risk using a 3-stage mathematical formula."),
    ("🔔", "Guardian Dashboard","Real-time risk alerts. Consent management. Timeline of risk events per adolescent."),
    ("🔒", "Privacy-First",     "Explicit opt-in consent. Encrypted transmission (TLS). No password or sensitive field capture."),
]
for i,(icon,title,body) in enumerate(sys_items):
    sy = Inches(1.55) + i*Inches(1.4)
    add_rect(sl, Inches(0.5), sy, Inches(5.8), Inches(1.3), C_CARD, C_PRIMARY, Pt(0.4))
    add_rect(sl, Inches(0.5), sy, Pt(4), Inches(1.3), C_PRIMARY)
    add_textbox(sl, icon,  Inches(0.65), sy+Inches(0.12), Inches(0.5), Inches(0.45), font_size=22, color=C_WHITE)
    add_textbox(sl, title, Inches(1.28), sy+Inches(0.12), Inches(4.9), Inches(0.28), font_size=11, bold=True, color=C_WHITE)
    add_textbox(sl, body,  Inches(1.28), sy+Inches(0.44), Inches(4.9), Inches(0.72), font_size=9.5, color=C_MUTED, wrap=True)

# Right: SUS
add_textbox(sl, "OBJECTIVE IV — SUS USABILITY RESULTS", Inches(6.6), Inches(1.22), Inches(6.2), Inches(0.26),
            font_size=8, bold=True, color=C_ACCENT)
add_rect(sl, Inches(6.6), Inches(1.55), Inches(6.2), Inches(3.45), C_CARD, C_ACCENT, Pt(0.8))

add_textbox(sl, "Average SUS Score", Inches(6.8), Inches(1.7), Inches(5.8), Inches(0.28),
            font_size=9, color=C_MUTED)
add_textbox(sl, "69.1", Inches(6.8), Inches(2.0), Inches(2.8), Inches(0.8),
            font_size=44, bold=True, color=C_ACCENT)
add_textbox(sl, "24 Participants", Inches(10.2), Inches(2.0), Inches(2.4), Inches(0.28),
            font_size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)

# Above Average badge
add_rect(sl, Inches(10.0), Inches(2.32), Inches(2.6), Inches(0.32),
         RGBColor(0x0D,0x25,0x1A), C_ACCENT, Pt(0.5))
add_textbox(sl, "✔ ABOVE AVERAGE", Inches(10.0), Inches(2.35), Inches(2.6), Inches(0.26),
            font_size=9, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Bar track
add_rect(sl, Inches(6.8), Inches(2.85), Inches(5.8), Inches(0.22), RGBColor(0x20,0x20,0x40), C_PRIMARY, Pt(0.2))
# Bar fill (69.1%)
add_rect(sl, Inches(6.8), Inches(2.85), Inches(5.8*0.691), Inches(0.22), C_ACCENT)
# Labels
sus_labels = [("0", Inches(6.8)), ("68 Avg", Inches(6.8) + Inches(5.8*0.68)), ("100", Inches(12.45))]
for lbl, lx in sus_labels:
    add_textbox(sl, lbl, lx-Inches(0.2), Inches(3.12), Inches(0.5), Inches(0.22),
                font_size=8, color=C_MUTED, align=PP_ALIGN.CENTER)

# User feedback note
add_rect(sl, Inches(6.8), Inches(3.42), Inches(5.8), Inches(1.38),
         RGBColor(0x14,0x14,0x35), C_PRIMARY, Pt(0.4))
add_textbox(sl,
    '💬 User Feedback:\n"Users want to see WHY a risk was assigned — not just the severity level."\n\n→ Identified as top recommendation for future work:\n   Explainable AI (XAI) Insights Panel',
    Inches(6.95), Inches(3.5), Inches(5.5), Inches(1.22),
    font_size=9.5, color=C_TEXT, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SIGNIFICANCE / CONTRIBUTION
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Significance", [("Contribution & ", C_WHITE), ("Impact", C_PRIMARY)])
add_textbox(sl, "10", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

add_textbox(sl, "FOR THE USER", Inches(0.5), Inches(1.22), Inches(5.8), Inches(0.26),
            font_size=8, bold=True, color=C_PRIMARY)
add_textbox(sl, "FOR THE FIELD", Inches(6.8), Inches(1.22), Inches(5.8), Inches(0.26),
            font_size=8, bold=True, color=C_ACCENT)

user_sigs = [
    ("👧", "Zero-Effort Protection",
     "Adolescents receive protection without any active engagement or self-reporting — removing the biggest barrier to mental health support."),
    ("👨‍👩‍👧", "Empowered Guardians",
     "Parents and carers receive timely, actionable alerts enabling earlier conversations and professional referrals before crises escalate."),
    ("🔒", "Privacy-First Design",
     "Explicit consent architecture ensures users retain full control over their data, building trust that existing monitoring tools lack."),
]

field_sigs = [
    ("🔬", "Novel RAG Application",
     "First application of FAISS-backed RAG retrieval to passive browser telemetry for adolescent mental health — a new Digital Phenotyping paradigm."),
    ("🏥", "Clinical Safety Framework",
     "Explicit ≥90% High-Risk Recall target aligned with clinical triage standards — bridging AI detection and clinical safety benchmarks."),
    ("🔭", "Explainability Roadmap",
     "User-driven recommendation for an XAI Insights Panel transforms the dashboard from an alerting tool to an educational clinical instrument."),
]

sh = Inches(1.62)
item_h = Inches(1.7)
for i,(icon,title,body) in enumerate(user_sigs):
    y = sh + i*(item_h+Inches(0.05))
    add_rect(sl, Inches(0.5), y, Inches(6.0), item_h, C_CARD, C_PRIMARY, Pt(0.4))
    add_textbox(sl, icon,  Inches(0.65), y+Inches(0.12), Inches(0.55), Inches(0.5), font_size=22, color=C_WHITE)
    add_textbox(sl, title, Inches(1.35), y+Inches(0.12), Inches(5.0), Inches(0.28), font_size=12, bold=True, color=C_WHITE)
    add_textbox(sl, body,  Inches(1.35), y+Inches(0.44), Inches(5.0), Inches(1.1), font_size=10, color=C_MUTED, wrap=True)

for i,(icon,title,body) in enumerate(field_sigs):
    y = sh + i*(item_h+Inches(0.05))
    add_rect(sl, Inches(6.8), y, Inches(6.0), item_h, C_CARD, C_ACCENT, Pt(0.4))
    add_textbox(sl, icon,  Inches(6.95), y+Inches(0.12), Inches(0.55), Inches(0.5), font_size=22, color=C_WHITE)
    add_textbox(sl, title, Inches(7.65), y+Inches(0.12), Inches(5.0), Inches(0.28), font_size=12, bold=True, color=C_WHITE)
    add_textbox(sl, body,  Inches(7.65), y+Inches(0.44), Inches(5.0), Inches(1.1), font_size=10, color=C_MUTED, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Conclusion", [("What Was ", C_WHITE), ("Accomplished", C_ACCENT)])
add_textbox(sl, "11", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

conclusions = [
    "Developed a fully functional browser extension that passively detects depressive language with 88.1% overall accuracy.",
    "Implemented a novel RAG-driven LLM pipeline (FAISS + Gemini), achieving a clinically critical High-Risk Recall of 94.1% — surpassing the ≥90% safety target.",
    "Built a complete guardian escalation framework with a real-time dashboard, consent management, and automated alert notifications.",
    "Validated usability with 24 participants achieving a SUS score of 69.1 (Above Average) — confirming accessibility for non-technical users.",
    "Identified key future work: an Explainable AI (XAI) Insights Panel to surface linguistic triggers — driven directly by user feedback.",
]

ct = Inches(1.28)
ch2 = Inches(0.88)
for i, text in enumerate(conclusions):
    y = ct + i*(ch2+Inches(0.06))
    add_rect(sl, Inches(0.5), y, Inches(12.3), ch2, C_CARD, C_PRIMARY, Pt(0.4))
    add_textbox(sl, "✔", Inches(0.65), y+Inches(0.2), Inches(0.35), Inches(0.45),
                font_size=18, bold=True, color=C_ACCENT)
    add_textbox(sl, text, Inches(1.15), y+Inches(0.14), Inches(11.5), ch2-Inches(0.22),
                font_size=11, color=C_TEXT, wrap=True)

# Summary banner
add_rect(sl, Inches(0.5), Inches(5.82), Inches(12.3), Inches(0.5),
         RGBColor(0x0D,0x25,0x1A), C_ACCENT, Pt(0.5))
add_textbox(sl,
    "🎯 MindGuard successfully addresses the critical gap in passive adolescent mental health monitoring — bridging AI detection and real-world guardian intervention.",
    Inches(0.65), Inches(5.9), Inches(12.0), Inches(0.38),
    font_size=10, bold=True, color=C_ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — REFERENCES
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "References", [("Works ", C_WHITE), ("Cited", C_PRIMARY), (" in This Presentation", C_WHITE)])
add_textbox(sl, "12", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

references = [
    "Q. Guo et al. (2024). Detection of suicidal ideation in social media text using BERT-based classification. Journal of Biomedical Informatics.",
    "Ilapaka, A., & Ghosh, S. (2025). RAG-enhanced conversational agents for mental health support using FAISS vector retrieval. Applied Intelligence.",
    "Zhang, Y. et al. (2023). SouLLMate: An LLM-based emotional support system for mental health applications. arXiv.",
    "Fitzpatrick, K. K. et al. (2022). Delivering cognitive behaviour therapy to young adults via a conversational agent, Woebot. JMIR Mental Health.",
    "Hollis, C. et al. (2021). Annual research review: Digital health interventions for children and young people with mental health problems. J. Child Psychology & Psychiatry, 62(4).",
    "Torous, J. et al. (2021). The growing field of digital psychiatry: apps, social media, chatbots, and virtual reality. World Psychiatry, 20(3), 318–335.",
    "O'Reilly, M. et al. (2023). Examining the risk and protective associations of social media use and adolescent mental health. npj Mental Health Research.",
    "Thabrew, H. et al. (2022). Digital health interventions for adolescent mental health. Child and Adolescent Mental Health, 27(2), 98–109.",
]

rt = Inches(1.25)
rh = Inches(0.72)
for i, ref in enumerate(references):
    y = rt + i*(rh+Inches(0.02))
    add_rect(sl, Inches(0.5), y, Pt(5), rh, C_PRIMARY)
    add_textbox(sl, ref, Inches(0.7), y+Inches(0.06), Inches(12.1), rh-Inches(0.1),
                font_size=10, color=C_MUTED, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — ACKNOWLEDGEMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
slide_header(sl, "Acknowledgement", [("With Sincere ", C_WHITE), ("Gratitude", C_PRIMARY)])
add_textbox(sl, "13", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

acks = [
    ("🙏", "Almighty God",         "For wisdom, strength, and the grace to complete this project."),
    ("👨‍🏫", "My Supervisor",       "For patient guidance, technical mentorship, and invaluable feedback throughout this research journey."),
    ("🏫", "Head of Department",   "For providing the academic structure and resources that enabled this project."),
    ("📚", "Faculty of Sciences",  "For the academic foundation and institutional support across four years of study."),
    ("👥", "Colleagues & Friends", "For the collaboration, encouragement, and moral support throughout the academic year."),
    ("❤️", "Family",               "For unwavering support, belief, and sacrifice that made this academic journey possible."),
]

aw = Inches(3.9)
ah = Inches(2.2)
for i,(icon,title,body) in enumerate(acks):
    col = i % 3
    row = i // 3
    ax = Inches(0.5) + col*(aw+Inches(0.15))
    ay = Inches(1.3) + row*(ah+Inches(0.18))
    add_rect(sl, ax, ay, aw, ah, C_CARD, C_PRIMARY, Pt(0.5))
    add_textbox(sl, icon,  ax, ay+Inches(0.2), aw, Inches(0.55), font_size=28, align=PP_ALIGN.CENTER, color=C_WHITE)
    add_textbox(sl, title, ax, ay+Inches(0.82), aw, Inches(0.32), font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, body,  ax+Inches(0.15), ay+Inches(1.18), aw-Inches(0.3), Inches(0.85),
                font_size=9.5, color=C_MUTED, align=PP_ALIGN.CENTER, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
add_textbox(sl, "14", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.3),
            font_size=10, color=C_PRIMARY, bold=True)

add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0x0A,0x0A,0x1E))
add_rect(sl, Inches(2), Inches(0.5), Inches(5), Inches(5), RGBColor(0x10,0x0A,0x30))

add_textbox(sl, "🧠🛡️", Inches(5.7), Inches(1.2), Inches(2), Inches(0.8),
            font_size=36, align=PP_ALIGN.CENTER, color=C_WHITE)

add_textbox(sl, "Thank You",
            Inches(1.5), Inches(2.1), Inches(10.3), Inches(1.2),
            font_size=56, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, "for your time and attention.",
            Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.4),
            font_size=16, color=C_MUTED, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(5.9), Inches(3.82), Inches(1.5), Pt(3), C_PRIMARY)

add_textbox(sl,
    "This research was conducted in partial fulfilment of the requirements\nfor the award of BSc Computer Science.",
    Inches(2.0), Inches(4.0), Inches(9.3), Inches(0.55),
    font_size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

chips = ["Adelekan Igbagboyemi Mary", "BSc Computer Science · 2026", "🛡️ MindGuard"]
chip_w = Inches(3.0)
for i, chip in enumerate(chips):
    cx = Inches(0.5) + (SLIDE_W - chip_w*3 - Inches(0.3)*2)/2 + i*(chip_w+Inches(0.3))
    add_rect(sl, cx, Inches(4.72), chip_w, Inches(0.38),
             RGBColor(0x1A,0x14,0x45), C_PRIMARY, Pt(0.5))
    add_textbox(sl, chip, cx, Inches(4.76), chip_w, Inches(0.3),
                font_size=11, bold=True, color=C_SECONDARY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = r"evaluation_form\MindGuard_Presentation.pptx"
prs.save(out_path)
print("Saved: " + out_path)
print("Total slides: " + str(len(prs.slides)))
