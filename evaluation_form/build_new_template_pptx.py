from pptx import Presentation

# Load the template
prs = Presentation('Final Project Slide to Use.pptx')

title_slide_layout = prs.slide_layouts[0]
content_layout = prs.slide_layouts[2] # Title and Content
two_content_layout = prs.slide_layouts[5] # Two Content

def add_slide(layout, title):
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    return slide

def set_bullets(tf, bullet_list):
    if not tf: return
    tf.text = bullet_list[0] if bullet_list else ''
    for b in bullet_list[1:]:
        p = tf.add_paragraph()
        p.text = b

slide = add_slide(title_slide_layout, 'MindGuard')
if len(slide.placeholders) > 1:
    slide.placeholders[1].text = 'Development of a Browser Extension for Detecting Depressive Language\nin Adolescents in Online Communication Using a RAG-Driven LLM\n\nAdelekan Igbagboyemi Mary\nComputer Science · Final Year · 2025/2026'

slide = add_slide(content_layout, 'Introduction: Background to the Field')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    'Global Mental Health Crisis: Depression affects 280 million people worldwide. Adolescents are disproportionately impacted.',
    'Adolescents & Online Communication: Adolescents spend 7+ hours/day online, frequently expressing emotional distress through written text.',
    'The Intervention Gap: Existing tools rely on active user engagement, but at-risk adolescents disengage when most distressed. Early detection requires passive monitoring.',
    'The digital footprint left by adolescents represents an untapped early-warning signal for identifying depressive risk before it escalates.'
])

slide = add_slide(content_layout, 'Literature Review: Key Works Reviewed')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    'Q. Guo et al. (2024): Suicidal ideation detection in social media. Gap: Not age-targeted; no guardian alerts. MindGuard: Age-targeted + Alerts.',
    'Ilapaka & Ghosh (2025): AI chatbot for depression support. Gap: Requires active user initiation. MindGuard: Passive monitoring.',
    'Zhang et al. (2023): LLM emotional support. Gap: No guardian escalation; overall F1=0.68. MindGuard: Guardian dashboard + 94.1% High-Risk F1.',
    'Fitzpatrick et al. (2022): CBT chatbot. Gap: Requires active engagement. MindGuard: Zero user effort required.',
    'Hollis et al. (2021): Digital mental health review. Gap: Lack of passive real-time browser tools. MindGuard: Browser extension solution.'
])

slide = add_slide(content_layout, 'Statement of Problem')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    'Active-Engagement Dependency: Existing tools work only if the user chooses to engage. Adolescents experiencing episodes are least likely to do so.',
    'Passive Expression Goes Undetected: No browser-based tool exists to passively capture and analyse real-time signals without user effort.',
    'No Guardian Escalation Framework: State-of-the-art detection models lack an integrated guardian notification system.',
    'Not Designed for Adolescents: Current tools ignore the informal, colloquial, and coded language unique to teenage online communication.'
])

slide = add_slide(content_layout, 'Aim')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    'To design, develop, and evaluate a browser extension that passively monitors online text for depressive language patterns in adolescents, using a RAG-driven Large Language Model, and automatically notifies designated guardians when high-risk indicators are detected.',
    'Key Features:',
    '  • Passive Monitoring',
    '  • RAG + LLM Integration',
    '  • Risk Classification',
    '  • Guardian Alerts'
])

slide = add_slide(content_layout, 'Objectives')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    '1. Collect & Pre-process Dataset: Dreaddit dataset (2,838 samples) of adolescent-style text.',
    '2. Design & Develop Risk Model: Hybrid NLP + RAG pipeline for accuracy and reliability.',
    '3. Implement System Architecture: Browser extension with secure guardian notification framework and consent management.',
    '4. Evaluate Usability: System Usability Scale (SUS) evaluation with target demographic.'
])

slide = add_slide(content_layout, 'Methodology')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    'Dataset: Dreaddit Reddit dataset. Text cleaning, stopword removal, FAISS embedding index.',
    'Risk Model: 3-class Risk Engine (Low / Moderate / High). FAISS vector search retrieves similar cases -> Gemini 1.5 Flash reasons over context (14-day temporal window).',
    'System: Browser extension (JS), FastAPI backend, Guardian dashboard (HTML/JS). Passive text interception -> encrypted API call -> conditional alert.',
    'Evaluation: Quantitative (Accuracy, Precision, Recall, F1) on 42-sample test set. Qualitative (SUS questionnaire with 24 participants).'
])

slide = add_slide(content_layout, 'Results: Model Performance')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    'Overall Metrics:',
    '  • Accuracy: 88.1%',
    '  • Precision: 87.4%',
    '  • Recall: 88.1%',
    '  • F1-Score: 87.5%',
    'Class Breakdown:',
    '  • Low Risk: F1 = 78.5%. 26.7% false-positive rate.',
    '  • Moderate Risk: F1 = 86.1%. Strong identification of borderline cases.',
    '  • High Risk: F1 = 94.1%, Recall = 94.1%. Meets >=90% clinical safety target.'
])

slide = add_slide(content_layout, 'Results: System & Usability')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    'System Delivered:',
    '  • Browser Extension (Passively captures text, encrypted)',
    '  • RAG Risk Engine (FAISS + Gemini LLM)',
    '  • Guardian Dashboard (Real-time alerts, consent management)',
    '  • Privacy-First Design (Explicit opt-in, no password capture)',
    'Usability Evaluation (SUS):',
    '  • Average SUS Score: 69.1 (Above Average) across 24 Participants.',
    '  • Key User Feedback: Users want an Explainable AI (XAI) Insights Panel to understand WHY a risk was assigned.'
])

slide = add_slide(two_content_layout, 'Significance: Contribution & Impact')
tf1 = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
tf2 = slide.placeholders[2].text_frame if len(slide.placeholders) > 2 else None
if tf1: set_bullets(tf1, [
    'For the User:',
    '  • Zero-Effort Protection: Adolescents receive protection without active engagement.',
    '  • Empowered Guardians: Timely, actionable alerts for earlier conversations.',
    '  • Privacy-First Design: Explicit consent builds trust.'
])
if tf2: set_bullets(tf2, [
    'For the Field:',
    '  • Novel RAG Application: First FAISS-backed RAG retrieval for passive telemetry.',
    '  • Clinical Safety Framework: Explicit >=90% High-Risk Recall target.',
    '  • Explainability Roadmap: XAI Insights Panel bridges alerting tool to educational instrument.'
])

slide = add_slide(content_layout, 'Conclusion')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    'Developed a fully functional browser extension detecting depressive language with 88.1% accuracy.',
    'Implemented novel RAG-driven LLM pipeline (FAISS + Gemini), achieving clinically critical High-Risk Recall of 94.1%.',
    'Built a complete guardian escalation framework with real-time dashboard and consent management.',
    'Validated usability (SUS score: 69.1), confirming accessibility.',
    'MindGuard successfully addresses the gap in passive adolescent mental health monitoring, bridging AI detection and real-world intervention.'
])

slide = add_slide(content_layout, 'References')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    'Q. Guo et al. (2024). Detection of suicidal ideation in social media text.',
    'Ilapaka, A., & Ghosh, S. (2025). RAG-enhanced conversational agents for mental health support.',
    'Zhang, Y. et al. (2023). SouLLMate: An LLM-based emotional support system.',
    'Fitzpatrick, K. K. et al. (2022). Delivering cognitive behaviour therapy via Woebot.',
    'Hollis, C. et al. (2021). Digital health interventions for children and young people.',
    'Torous, J. et al. (2021). The growing field of digital psychiatry.'
])

slide = add_slide(content_layout, 'Acknowledgement')
tf = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
if tf: set_bullets(tf, [
    'Almighty God: For wisdom, strength, and grace.',
    'My Supervisor: For patient guidance and technical mentorship.',
    'Head of Department: For providing the academic structure.',
    'Faculty of Sciences: For the academic foundation.',
    'Colleagues & Friends: For collaboration and moral support.',
    'Family: For unwavering support and belief.'
])

slide = add_slide(title_slide_layout, 'Thank You')
if len(slide.placeholders) > 1:
    slide.placeholders[1].text = 'Questions & Answers'

prs.save('MindGuard_Presentation_Template_Applied.pptx')
